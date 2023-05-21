""" TODO
1. Take the path of a .pptx file.          --> DONE
2. Parse the presentation to get its data. --> DONE
3. Go through every slide separately, and:
4.          Extract all the text from all text boxes.
5.          Insert the text into an appropriate prompt for GPT.
6.          Send the prompt in a request to the OpenAI API.
7.          Extract the AI's reply from the response.
8. Gather the explanations for all the slides in a list.
9. Save the list in a JSON file.
"""


def get_file_path():
    import os

    while True:
        file_path = input("Enter the path to the presentation file: ")
        if file_path.endswith(".pptx") and os.path.exists(file_path):
            return file_path
        print("Invalid path/ file not found, please try again")


def extract_text_from_presentation(file_path):
    from pptx import Presentation

    # ValueError if file_path is not a valid pptx file, validation is done in get_file_path()
    prs = Presentation(file_path)
    prs_content = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text.strip())

        prs_content.append(','.join(slide_text))
    return prs_content


def ai_text_analyzer(presentation):
    import openai
    import time

    OPENAI_API_KEY = "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"

    openai.api_key = OPENAI_API_KEY
    system_prompt = "You're an AI text analyzer assisting with presentation summarization.For each slide's content you"\
                    " receive, generate a concise summary of the text.\n User:'Slide content'\nAI:'Summary explanation'"
    messages = [{"role": "system", "content": system_prompt}]  # the behavior of the system
    responses = []

    for index, slide_content in enumerate(presentation):
        # set instructions to the chat as a user message (the slide content)
        messages.append({"role": "user", "content": slide_content})
        # sends the request
        completion = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages)
        # get response
        chat_response = completion.choices[0].message.content
        # keeping the history of the chat
        messages.append({"role": "assistant", "content": chat_response})
        responses.append(chat_response + "\n")
        time.sleep(60)

    return responses


def main():
    print("-" * 20, "Hello ChatBot", "-" * 20)
    file_path = get_file_path()
    presentation_content = extract_text_from_presentation(file_path)
    print(ai_text_analyzer(presentation_content))


if __name__ == "__main__":
    main()
