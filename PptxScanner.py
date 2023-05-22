from pptx import Presentation


class PptxScanner:

    def __init__(self, path):
        self._presentation_path = path
        self._presentation = self.open_presentation(path)
        self._prs_content = []

    @staticmethod
    def open_presentation(path):
        """ Opens a presentation file and returns a Presentation object"""
        import os
        if path.endswith(".pptx") and os.path.exists(path):
            return Presentation(path)
        else:
            raise Exception("File path is not a valid .pptx type")

    def scan_presentation(self):
        """ Scans the presentation and collect the text from each slide"""
        self._prs_content = [self._get_slide_content(slide) for slide in self._presentation.slides]
        return self._prs_content

    @staticmethod
    def _get_slide_content(slide):
        """ Returns the text from a slide"""
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text.strip())
        return ''.join(slide_text)

    @property
    def get_path(self):
        return self._presentation_path
