"""
        PptxScanner.py
    ------------------------
    The file that responsible for
    scanning the presentation.
"""
from pptx import Presentation


class PptxScanner:
    def __init__(self, path: str):
        self._presentation_path = path
        self._presentation = self.open_presentation(path)
        self._prs_content = []

    @staticmethod
    def open_presentation(path: str) -> Presentation or Exception():
        """ Opens a presentation file and returns a Presentation object
        :param path: the path of the presentation
        :return: Presentation object
        """
        import os
        if path.endswith(".pptx") and os.path.exists(path):
            return Presentation(path)
        else:
            raise Exception("File path is not a valid .pptx type")

    def scan_presentation(self):
        """ Scans the presentation and collect the text from each slide"""
        self._prs_content = [self._get_slide_content(slide) for slide in self._presentation.slides]
        return self._prs_content

    @staticmethod
    def _get_slide_content(slide: Presentation) -> str:
        """ Returns the text from a slide
        :param slide: a slide object
        :return: the text from the slide
        """
        slide_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text.append(run.text.strip())
        return ''.join(slide_text)

    @property
    def get_path(self) -> str:
        return self._presentation_path
